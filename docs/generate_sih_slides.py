"""
Script to generate high-quality Smart India Hackathon 2026 Technical Approach PowerPoint Slides
for ISRO Problem Statement ID: 26169 (AI-Based Virtual Camera Tracking System for Mobile FSOC Terminals).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_sih_deck():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (Dark Theme / Aerospace Cyan & Emerald)
    COLOR_BG = RGBColor(7, 10, 20)          # #070a14
    COLOR_CARD_BG = RGBColor(14, 22, 40)     # #0e1628
    COLOR_CARD_BORDER = RGBColor(56, 189, 248) # #38bdf8
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_CYAN = RGBColor(56, 189, 248)
    COLOR_EMERALD = RGBColor(52, 211, 153)
    COLOR_AMBER = RGBColor(251, 191, 36)
    COLOR_MUTED = RGBColor(148, 163, 184)
    COLOR_DARK_BOX = RGBColor(20, 30, 55)

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, subtitle_text="SMART INDIA HACKATHON 2026 | ISRO PS ID: 26169"):
        # Header background line / container
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.95), Inches(12.533), Inches(0.02)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = COLOR_CARD_BORDER
        header_line.line.fill.background()

        # Team Badge
        team_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.25), Inches(2.2), Inches(0.55)
        )
        team_box.fill.solid()
        team_box.fill.fore_color.rgb = COLOR_DARK_BOX
        team_box.line.color.rgb = COLOR_CYAN
        team_box.line.width = Pt(1.5)
        tf = team_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "🚀 TEAM FALCONS"
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(Inches(2.8), Inches(0.15), Inches(7.7), Inches(0.75))
        tf_title = title_box.text_frame
        p_t = tf_title.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = "Arial"
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.alignment = PP_ALIGN.CENTER

        p_sub = tf_title.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_EMERALD
        p_sub.alignment = PP_ALIGN.CENTER

        # ISRO Badge
        isro_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.733), Inches(0.25), Inches(2.2), Inches(0.55)
        )
        isro_box.fill.solid()
        isro_box.fill.fore_color.rgb = COLOR_DARK_BOX
        isro_box.line.color.rgb = COLOR_AMBER
        isro_box.line.width = Pt(1.5)
        tf_isro = isro_box.text_frame
        p_isro = tf_isro.paragraphs[0]
        p_isro.text = "🇮🇳 ISRO / DOS"
        p_isro.font.name = "Arial"
        p_isro.font.size = Pt(12)
        p_isro.font.bold = True
        p_isro.font.color.rgb = COLOR_AMBER
        p_isro.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 1: TECHNICAL APPROACH (Tech Stack + System Architecture Data Flow)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    add_header(slide1, "TECHNICAL APPROACH")

    # LEFT COLUMN: Tech Stack Title
    left_title = slide1.shapes.add_textbox(Inches(0.4), Inches(1.05), Inches(5.8), Inches(0.4))
    tf_lt = left_title.text_frame
    p_lt = tf_lt.paragraphs[0]
    p_lt.text = "❖ TECHNOLOGY STACK"
    p_lt.font.name = "Arial"
    p_lt.font.size = Pt(15)
    p_lt.font.bold = True
    p_lt.font.color.rgb = COLOR_CYAN

    # Tech Stack Cards Grid
    tech_items = [
        ("🌐 Frontend & 3D Twin", "Three.js (WebGL), HTML5 Canvas HUD, Chart.js, Vanilla CSS Glassmorphism, IBM Plex typography, 60 FPS WebSockets.", 1.45, 0.4, 2.85, 1.1),
        ("⚡ Backend Engine", "Python 3.11+, FastAPI (ASGI), Uvicorn Server, Asyncio Event Loop, NumPy, SciPy (Vectorized Linear Algebra).", 1.45, 3.35, 2.85, 1.1),
        ("🧠 AI & Perception", "YOLOv8-FSOC Beacon Detector, ONNX Runtime, TensorRT edge acceleration, Pinhole Camera Intrinsic Model.", 2.65, 0.4, 2.85, 1.1),
        ("🎯 Estimation & Control", "6-DOF Extended Kalman Filter (EKF), Dual-Axis PID Controller with Anti-Windup & Slew Clamping (<= 45 deg/s).", 2.65, 3.35, 2.85, 1.1),
        ("🔬 Optical Physics & Link Budget", "1550nm C-Band EDFA Model, Beer-Lambert Extinction, Gaussian Beam Profile, SNR and Complementary Error Function BER (Pe = 1/2 erfc(sqrt(SNR)/2sqrt(2))).", 3.85, 0.4, 5.8, 1.05),
        ("🔌 Hardware & Benchmarks", "STM32 / ESP32 Micro-ROS, Serial CAN-Bus / gRPC, Pan-Tilt Brushless Gimbals, Pytest Suite, Automated 1-Click ISRO Benchmark Logger.", 5.00, 0.4, 5.8, 0.95),
    ]

    for title, desc, top, left, width, height in tech_items:
        card = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_DARK_BOX
        card.line.width = Pt(1)
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.1)
        tf_c.margin_right = Inches(0.1)
        tf_c.margin_top = Inches(0.08)
        tf_c.margin_bottom = Inches(0.08)

        p_head = tf_c.paragraphs[0]
        p_head.text = title
        p_head.font.name = "Arial"
        p_head.font.size = Pt(10.5)
        p_head.font.bold = True
        p_head.font.color.rgb = COLOR_EMERALD

        p_body = tf_c.add_paragraph()
        p_body.text = desc
        p_body.font.name = "Arial"
        p_body.font.size = Pt(9.5)
        p_body.font.color.rgb = COLOR_WHITE

    # Callout Box at bottom left
    summary_box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(6.05), Inches(5.8), Inches(1.15)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = COLOR_DARK_BOX
    summary_box.line.color.rgb = COLOR_AMBER
    summary_box.line.width = Pt(1.5)
    tf_s = summary_box.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = Inches(0.12)
    tf_s.margin_top = Inches(0.08)
    p_s = tf_s.paragraphs[0]
    p_s.text = "❖ FSOC AI Virtual Camera Tracking Engine Summary:"
    p_s.font.name = "Arial"
    p_s.font.size = Pt(10)
    p_s.font.bold = True
    p_s.font.color.rgb = COLOR_AMBER

    p_sb = tf_s.add_paragraph()
    p_sb.text = "Custom YOLOv8 for sub-pixel beacon localization; 6-DOF EKF trajectory prediction during 2.0s cloud occlusions; dual-axis PID closed-loop actuation at 60 FPS (< 5ms latency) maintaining pointing error <= 8.72 mrad."
    p_sb.font.name = "Arial"
    p_sb.font.size = Pt(9)
    p_sb.font.color.rgb = COLOR_WHITE

    # RIGHT COLUMN: Closed-Loop System Architecture
    right_title = slide1.shapes.add_textbox(Inches(6.6), Inches(1.05), Inches(6.3), Inches(0.4))
    tf_rt = right_title.text_frame
    p_rt = tf_rt.paragraphs[0]
    p_rt.text = "❖ CLOSED-LOOP SYSTEM ARCHITECTURE & DATA FLOW"
    p_rt.font.name = "Arial"
    p_rt.font.size = Pt(15)
    p_rt.font.bold = True
    p_rt.font.color.rgb = COLOR_CYAN

    # Architecture Pipeline Flow Steps
    steps = [
        ("01", "3D Kinematics & Channel Engine", "Target trajectory [X,Y,Z] in world frame, weather extinction, platform jitter", COLOR_CYAN),
        ("02", "Virtual Pinhole Camera Sensor", "Projects 3D coordinates to 2D image plane (1920x1080) using intrinsic matrix K", COLOR_CYAN),
        ("03", "AI Beacon Detector (YOLOv8-FSOC)", "Sub-pixel centroid & bounding box extraction (< 2.1 ms inference latency)", COLOR_EMERALD),
        ("04", "6-DOF Extended Kalman Filter (EKF)", "State estimation [u, v, u_dot, v_dot, u_ddot, v_ddot] + 2.0s occlusion extrapolation", COLOR_EMERALD),
        ("05", "Dual-Axis PID Gimbal Controller", "Clamped slew (<= 45 deg/s) to maintain optical alignment within 8.72 mrad (0.50 deg)", COLOR_AMBER),
        ("06", "1550nm FSOC Optical Physics & Link", "Evaluates Beer-Lambert loss, RSSI (dBm), SNR (dB), and validates BER <= 1e-9", COLOR_CYAN),
    ]

    top_pos = 1.45
    for num, title, desc, col in steps:
        step_box = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(top_pos), Inches(6.3), Inches(0.68)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = COLOR_CARD_BG
        step_box.line.color.rgb = col
        step_box.line.width = Pt(1.5)
        tf_step = step_box.text_frame
        tf_step.word_wrap = True
        tf_step.margin_left = Inches(0.12)
        tf_step.margin_top = Inches(0.06)

        p_st = tf_step.paragraphs[0]
        p_st.text = f"[{num}]  {title}"
        p_st.font.name = "Arial"
        p_st.font.size = Pt(10.5)
        p_st.font.bold = True
        p_st.font.color.rgb = col

        p_sd = tf_step.add_paragraph()
        p_sd.text = desc
        p_sd.font.name = "Arial"
        p_sd.font.size = Pt(9)
        p_sd.font.color.rgb = COLOR_WHITE

        top_pos += 0.76

    # Bottom KPI Cards Row
    kpis = [
        ("≤ 8.72 mrad", "Coarse Threshold"),
        ("60 FPS", "Telemetry Rate"),
        ("< 5.0 ms", "Processing Latency"),
        ("> 95.0%", "Lock Retention"),
    ]

    kpi_left = 6.6
    for val, lbl in kpis:
        kbox = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(kpi_left), Inches(6.15), Inches(1.48), Inches(1.05)
        )
        kbox.fill.solid()
        kbox.fill.fore_color.rgb = COLOR_DARK_BOX
        kbox.line.color.rgb = COLOR_CARD_BORDER
        kbox.line.width = Pt(1)
        tf_k = kbox.text_frame
        tf_k.word_wrap = True
        tf_k.margin_top = Inches(0.15)
        p_kv = tf_k.paragraphs[0]
        p_kv.text = val
        p_kv.font.name = "Arial"
        p_kv.font.size = Pt(14)
        p_kv.font.bold = True
        p_kv.font.color.rgb = COLOR_EMERALD
        p_kv.alignment = PP_ALIGN.CENTER

        p_kl = tf_k.add_paragraph()
        p_kl.text = lbl
        p_kl.font.name = "Arial"
        p_kl.font.size = Pt(8.5)
        p_kl.font.color.rgb = COLOR_MUTED
        p_kl.alignment = PP_ALIGN.CENTER

        kpi_left += 1.6

    # =========================================================================
    # SLIDE 2: TECHNICAL APPROACH & METHODOLOGY (Matching Reference Image 2)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "TECHNICAL APPROACH & METHODOLOGY")

    # 3-Column Layout: Left (Methodology), Middle (Cyber-Physical), Right (Tech Cards)
    col1_w = Inches(3.9)
    col2_w = Inches(4.3)
    col3_w = Inches(4.0)

    # Column 1: Methodology Flow
    c1_title = slide2.shapes.add_textbox(Inches(0.4), Inches(1.05), col1_w, Inches(0.4))
    p_c1 = c1_title.text_frame.paragraphs[0]
    p_c1.text = "❖ 6-PHASE METHODOLOGY"
    p_c1.font.name = "Arial"
    p_c1.font.size = Pt(14)
    p_c1.font.bold = True
    p_c1.font.color.rgb = COLOR_CYAN

    method_steps = [
        ("1. Kinematics & Channel", "3D target trajectory + wind turbulence & weather attenuation."),
        ("2. Camera Projection", "Pinhole intrinsic projection matrix mapping [X,Y,Z] -> [u,v]."),
        ("3. AI YOLOv8 Detection", "Sub-pixel centroid localization with optical beacon filter."),
        ("4. 6-DOF EKF Estimation", "Forward trajectory extrapolation across 2.0s cloud occlusions."),
        ("5. Dual-Axis PID Control", "Clamped slew (<= 45 deg/s) to achieve pointing error <= 8.72 mrad."),
        ("6. Optical Link Validation", "1550nm link budget: RSSI, SNR, and BER <= 1e-9 (10 Gbps Lock)."),
    ]

    m_top = 1.45
    for title, desc in method_steps:
        mbox = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(m_top), col1_w, Inches(0.85)
        )
        mbox.fill.solid()
        mbox.fill.fore_color.rgb = COLOR_CARD_BG
        mbox.line.color.rgb = COLOR_DARK_BOX
        mbox.line.width = Pt(1)
        tf_m = mbox.text_frame
        tf_m.word_wrap = True
        tf_m.margin_left = Inches(0.1)
        tf_m.margin_top = Inches(0.06)

        p_mt = tf_m.paragraphs[0]
        p_mt.text = title
        p_mt.font.name = "Arial"
        p_mt.font.size = Pt(10)
        p_mt.font.bold = True
        p_mt.font.color.rgb = COLOR_EMERALD

        p_md = tf_m.add_paragraph()
        p_md.text = desc
        p_md.font.name = "Arial"
        p_md.font.size = Pt(8.5)
        p_md.font.color.rgb = COLOR_WHITE
        m_top += 0.92

    # Column 2: Cyber-Physical Mapping (Simulation vs Hardware)
    c2_title = slide2.shapes.add_textbox(Inches(4.6), Inches(1.05), col2_w, Inches(0.4))
    p_c2 = c2_title.text_frame.paragraphs[0]
    p_c2.text = "❖ SIMULATION TO HARDWARE MAPPING"
    p_c2.font.name = "Arial"
    p_c2.font.size = Pt(14)
    p_c2.font.bold = True
    p_c2.font.color.rgb = COLOR_CYAN

    hw_mappings = [
        ("Optical Perception", "Synthetic pinhole projection", "1550nm InGaAs SWIR Coarse Acquisition Camera", COLOR_CYAN),
        ("AI Target Detection", "YOLOv8-FSOC ONNX model", "Jetson Orin / Embedded FPGA AI Accelerator", COLOR_EMERALD),
        ("State Estimation", "6-DOF EKF Riccati filter", "Real-time C++ EKF flight controller firmware", COLOR_EMERALD),
        ("Gimbal Pointing", "Dual-axis PID with slew limit", "Direct CAN-bus PWM to Pan-Tilt Direct Motors", COLOR_AMBER),
        ("Optical Physics", "Gaussian beam & Beer-Lambert", "100mW EDFA Laser Diode + APD Photodiode", COLOR_CYAN),
        ("Mission Telemetry", "60 FPS WebSocket data bus", "ISTRAC Ground Station CCSDS / ZeroMQ Protocol", COLOR_MUTED),
    ]

    h_top = 1.45
    for sub, sim, hw, col in hw_mappings:
        hbox = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(h_top), col2_w, Inches(0.85)
        )
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = COLOR_CARD_BG
        hbox.line.color.rgb = col
        hbox.line.width = Pt(1)
        tf_h = hbox.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = Inches(0.1)
        tf_h.margin_top = Inches(0.06)

        p_ht = tf_h.paragraphs[0]
        p_ht.text = f"{sub}:"
        p_ht.font.name = "Arial"
        p_ht.font.size = Pt(10)
        p_ht.font.bold = True
        p_ht.font.color.rgb = col

        p_hd = tf_h.add_paragraph()
        p_hd.text = f"• Sim: {sim}\n• Hardware: {hw}"
        p_hd.font.name = "Arial"
        p_hd.font.size = Pt(8.5)
        p_hd.font.color.rgb = COLOR_WHITE
        h_top += 0.92

    # Column 3: Technologies Used
    c3_title = slide2.shapes.add_textbox(Inches(9.1), Inches(1.05), col3_w, Inches(0.4))
    p_c3 = c3_title.text_frame.paragraphs[0]
    p_c3.text = "❖ TECHNOLOGIES USED"
    p_c3.font.name = "Arial"
    p_c3.font.size = Pt(14)
    p_c3.font.bold = True
    p_c3.font.color.rgb = COLOR_CYAN

    tech_cards = [
        ("🎨 FrontEnd", "React / Three.js (WebGL), HTML5 Canvas Telescope HUD, Chart.js, Vanilla CSS Glassmorphism.", COLOR_CYAN),
        ("⚡ BackEnd", "Python 3.11+, FastAPI (ASGI), Uvicorn, NumPy, SciPy (Vectorized Math), WebSockets.", COLOR_EMERALD),
        ("🧠 AI & Estimation", "YOLOv8-FSOC, 6-DOF Extended Kalman Filter (EKF), Dual-Axis PID Slew Controller.", COLOR_EMERALD),
        ("📊 Benchmarking", "Automated 1-Click ISRO Performance Logger (JSON Summary, CSV Dataset, PDF Report).", COLOR_AMBER),
        ("📡 Optical Physics", "1550nm FSOC Laser Link Budget, Beer-Lambert Extinction, Gaussian Beam Loss, SNR & BER.", COLOR_CYAN),
    ]

    t_top = 1.45
    for title, desc, col in tech_cards:
        tbox = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(t_top), col3_w, Inches(1.03)
        )
        tbox.fill.solid()
        tbox.fill.fore_color.rgb = COLOR_DARK_BOX
        tbox.line.color.rgb = col
        tbox.line.width = Pt(1.2)
        tf_t = tbox.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = Inches(0.1)
        tf_t.margin_top = Inches(0.08)

        p_tt = tf_t.paragraphs[0]
        p_tt.text = title
        p_tt.font.name = "Arial"
        p_tt.font.size = Pt(10.5)
        p_tt.font.bold = True
        p_tt.font.color.rgb = col

        p_td = tf_t.add_paragraph()
        p_td.text = desc
        p_td.font.name = "Arial"
        p_td.font.size = Pt(9)
        p_td.font.color.rgb = COLOR_WHITE
        t_top += 1.10

    # =========================================================================
    # SLIDE 3: ALGORITHMIC RIGOR & MATHEMATICAL FORMULATION
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "ALGORITHMIC RIGOR & MATHEMATICAL FORMULATION", "MATHEMATICAL RIGOR & CONTROL EQUATIONS | ISRO PS ID: 26169")

    # 4 Algorithm Cards Grid (2x2)
    cards_s3 = [
        ("1. Pinhole Camera Projection Model", 
         "• 3D to 2D Homogeneous Coordinate Transform:\n"
         "  [u, v, 1]^T = (1/Z_c) * K * [R | T] * [X_w, Y_w, Z_w, 1]^T\n"
         "• Intrinsic Matrix K = diag(f_x, f_y, 1) + [0, 0, c_x; 0, 0, c_y; 0, 0, 0]\n"
         "• Optical Sensor Noise: eta ~ N(0, sigma_pixel^2) with sigma = 0.5 px\n"
         "• Real-time sub-pixel centroid extraction at 1920x1080 resolution.",
         1.2, 0.4, 6.1, 2.7, COLOR_CYAN),
        
        ("2. 6-DOF Extended Kalman Filter (EKF)", 
         "• State Vector: x = [u, v, u_dot, v_dot, u_ddot, v_ddot]^T\n"
         "• State Transition: x_{k|k-1} = F * x_{k-1} + w_k,  w_k ~ N(0, Q)\n"
         "• Measurement Update: z_k = H * x_k + v_k,  v_k ~ N(0, R)\n"
         "• Kalman Gain: K_k = P_{k|k-1} H^T (H P_{k|k-1} H^T + R)^{-1}\n"
         "• 2.0s Occlusion Extrapolation: P_k updates via continuous Riccati propagation.",
         1.2, 6.8, 6.1, 2.7, COLOR_EMERALD),
        
        ("3. Dual-Axis PID Gimbal Slew Dynamics", 
         "• Pixel Error to Angular Rate: e_theta = (u - c_x)/f_x,  e_phi = (v - c_y)/f_y\n"
         "• Control Law: u(t) = K_p e(t) + K_i int_0^t e(tau) dtau + K_d (de/dt)\n"
         "• Dynamic Anti-Windup: Clamps integral accumulator during high slews\n"
         "• Slew Rate Limiter: |omega_az| <= 45 deg/s, |omega_el| <= 45 deg/s\n"
         "• Pointing Error Target: theta_pointing <= 8.72 mrad (0.50 deg).",
         4.1, 0.4, 6.1, 2.7, COLOR_AMBER),
        
        ("4. 1550nm FSOC Optical Physics & Link Budget", 
         "• Received Power: P_RX(dBm) = P_TX + G_TX + G_RX - FSPL - L_pointing - A_atm\n"
         "• Gaussian Pointing Loss: L_pointing = -8.686 * (theta_pointing / theta_div)^2 dB\n"
         "• Beer-Lambert Extinction: A_atm(R) = 10 * gamma_ext * R * log10(e)\n"
         "• Bit Error Rate (BER): P_e = 0.5 * erfc(sqrt(SNR) / (2 * sqrt(2)))\n"
         "• Carrier Lock: BER <= 1e-9 & SNR >= 25 dB enables 10 Gbps data throughput.",
         4.1, 6.8, 6.1, 2.7, COLOR_CYAN),
    ]

    for title, body, top, left, width, height, col in cards_s3:
        box = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(11.5)
        p_t.font.bold = True
        p_t.font.color.rgb = col

        p_b = tf.add_paragraph()
        p_b.text = body
        p_b.font.name = "Courier New"
        p_b.font.size = Pt(9)
        p_b.font.color.rgb = COLOR_WHITE

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "SIH2026_Technical_Approach_ISRO_PS26169.pptx")
    prs.save(output_path)
    print(f"PowerPoint Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    create_sih_deck()

